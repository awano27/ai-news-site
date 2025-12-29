from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "2025年 AIの進化と歴史"
    subtitle.text = "主要なマイルストーンとブレイクスルー\n作成日: 2025年12月23日"

    # Helper function to add a bullet slide
    def add_bullet_slide(heading, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = heading
        tf = body_shape.text_frame
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(24)

    # Helper function to add a timeline slide
    def add_timeline_slide(heading, events):
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        shapes.title.text = heading
        
        # Draw a horizontal line
        left = Inches(1)
        top = Inches(3.5)
        width = Inches(8)
        height = Inches(0.1)
        shape = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192) # Blue
        
        # Add events
        for i, (year, desc) in enumerate(events):
            x = Inches(1.5 + i * 2)
            y_year = Inches(2.5)
            y_desc = Inches(4.0)
            
            # Year bubble
            oval = shapes.add_shape(MSO_SHAPE.OVAL, x, y_year, Inches(1), Inches(0.8))
            oval.fill.solid()
            oval.fill.fore_color.rgb = RGBColor(255, 192, 0) # Orange
            oval.text = year
            
            # Description box
            textbox = shapes.add_textbox(x - Inches(0.5), y_desc, Inches(2), Inches(1.5))
            tf = textbox.text_frame
            tf.text = desc
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Helper function to add a chart slide
    def add_chart_slide(heading, chart_data_dict):
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = heading
        
        chart_data = CategoryChartData()
        chart_data.categories = list(chart_data_dict.keys())
        chart_data.add_series('Adoption Rate', list(chart_data_dict.values()))
        
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

    # Helper function to add a diagram slide (Multimodal)
    def add_diagram_slide(heading):
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = heading
        
        # Central Hub
        center_x, center_y = Inches(4), Inches(3)
        radius = Inches(1.5)
        hub = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, center_x, center_y, radius*1.5, radius*1.5)
        hub.text = "Multimodal AI\nModel"
        hub.fill.solid()
        hub.fill.fore_color.rgb = RGBColor(112, 48, 160) # Purple
        
        # Inputs
        inputs = [("Text", Inches(4), Inches(1)), 
                  ("Image", Inches(2), Inches(4.5)), 
                  ("Audio", Inches(6), Inches(4.5))]
        
        for text, x, y in inputs:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.5), Inches(1))
            box.text = text
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0, 176, 80) # Green
            
            # Connector (simple line)
            # Note: Connectors are tricky in simple script, just placing shapes close is enough for "diagram" feel or drawing lines
            # Let's draw a line from box center to hub center
            # Simple approach: just shapes for now as requested "visual grouping"

    # Slide 2: Overview
    add_bullet_slide("2025年の概要", [
        "自律型AIエージェントの台頭：チャットから行動へ",
        "マルチモーダルAIの進化：視覚・聴覚・言語の統合",
        "産業への深い統合：医療、金融、開発現場での実用化",
        "AIの最適化と民主化：より効率的で身近な存在へ"
    ])

    # Slide 3: Timeline (Visual)
    add_timeline_slide("2025年の進化タイムライン", [
        ("Q1", "マルチモーダル\nモデルの普及"),
        ("Q2", "自律エージェント\n実用化開始"),
        ("Q3", "産業特化型\nAIの拡大"),
        ("Q4", "AGIに向けた\n新たな議論")
    ])

    # Slide 4: Chart (Visual)
    add_chart_slide("AIエージェントの企業採用率予測", {
        "2023": 5,
        "2024": 10,
        "2025": 25,
        "2026": 40
    })

    # Slide 5: Diagram (Visual)
    add_diagram_slide("マルチモーダルAIの統合イメージ")

    # Slide 6: 自律型AIエージェント (Text)
    add_bullet_slide("自律型AIエージェントの詳細", [
        "「Agentic AI」へのシフト：指示待ちから自律的な意思決定へ",
        "複雑なタスクの実行：スケジューリング、ソフトウェア開発、交渉など",
        "企業の採用拡大：Deloitteの予測では2025年に採用率が25%に到達"
    ])

    # Slide 7: 産業分野での革新 (Text)
    add_bullet_slide("産業分野での革新", [
        "医療：膨大な臨床データ解析による早期診断と個別化医療",
        "金融：自律型金融システムによる予測と高度な不正検知",
        "ソフトウェア開発：AIがコード生成・テスト・修正を自律的に実行"
    ])

    # Slide 8: まとめ
    add_bullet_slide("まとめ：2025年の位置づけ", [
        "AIは「ツール」から自律的な「パートナー」へと進化",
        "実験段階から実用・最適化のフェーズへ移行",
        "B2B AIの民主化により、あらゆる規模の企業で活用が進む",
        "次世代（AGI）に向けた重要な転換点"
    ])

    output_path = "d:\\ai-news-site-main\\AI_Evolution_2025.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
